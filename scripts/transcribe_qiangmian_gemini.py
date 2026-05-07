#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
宗教史讀書會 Gemini 音訊轉錄流水線 v2

使用 Gemini Audio 取代 Whisper，搭配 PPT 投影片作為轉錄參考。
轉錄品質更高，繁體中文更準確，段落結構更清晰。
Upsert 取代舊資料。

用法：
  python scripts/transcribe_qiangmian_gemini.py --list
  python scripts/transcribe_qiangmian_gemini.py --episode 1
  python scripts/transcribe_qiangmian_gemini.py --episode 1-5
  python scripts/transcribe_qiangmian_gemini.py --all
"""

import argparse
import re
import shutil
import subprocess
import sys
import tempfile
import time
from pathlib import Path

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

# ── 路徑設定 ──────────────────────────────────────────────────
REPO_DIR     = Path(__file__).parent.parent
PPT_DIR      = Path("G:/我的雲端硬碟/創作/千面上帝/宗教史讀書會")
TMP_DIR      = REPO_DIR / "_tmp_audio" / "qiangmian"
TMP_DIR.mkdir(parents=True, exist_ok=True)

PLAYLIST_URL = "https://www.youtube.com/playlist?list=PLNdU3g_-OSshfnyOakO5exMMvnSNeuIjZ"

# ── .env ──────────────────────────────────────────────────────
_env: dict[str, str] = {}
for _line in (REPO_DIR / ".env").read_text(encoding="utf-8").splitlines():
    if "=" in _line and not _line.startswith("#"):
        _k, _, _v = _line.partition("=")
        _env[_k.strip()] = _v.strip().strip('"')

SUPABASE_URL     = _env["SUPABASE_URL"]
SERVICE_ROLE_KEY = _env["SUPABASE_SERVICE_ROLE_KEY"]


# ── Gemini keys ───────────────────────────────────────────────
def _find_keys() -> list[str]:
    raw: list[str] = []
    for name in ("GEMINI_API_KEY", "Gemini_API_Key"):
        v = _env.get(name)
        if v:
            raw.append(v); break
    for n in range(1, 11):
        for base in ("Gemini_API_Key", "GEMINI_API_KEY"):
            v = _env.get(f"{base}_{n}")
            if v:
                raw.append(v); break
    keys: list[str] = []
    seen: set[str] = set()
    for r in raw:
        for piece in r.split(","):
            k = piece.strip()
            if k and k not in seen:
                seen.add(k); keys.append(k)
    return keys

GEMINI_KEYS = _find_keys()


# ── PPT ───────────────────────────────────────────────────────
def sorted_ppt_files() -> list[Path]:
    files = list(PPT_DIR.glob("*.pptx"))
    files.sort(key=lambda p: p.name[:10])  # sort by YYYY.MM.DD prefix
    return files


def load_ppt_text(ppt_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
        from pptx import Presentation
    prs = Presentation(str(ppt_path))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides_text.append(f"[投影片{i}] " + " / ".join(parts))
    return "\n".join(slides_text)


# ── YouTube ───────────────────────────────────────────────────
def fetch_playlist() -> list[dict]:
    import yt_dlp
    with yt_dlp.YoutubeDL({"quiet": True, "extract_flat": True}) as ydl:
        info = ydl.extract_info(PLAYLIST_URL, download=False)
    return info.get("entries", [])


def safe_name(s: str) -> str:
    return re.sub(r'[\\/:*?"<>|]', '_', s)[:80].strip()


def download_audio(video_id: str, label: str) -> Path | None:
    import yt_dlp
    for ext in ["m4a", "webm", "opus", "mp4"]:
        p = TMP_DIR / f"{label}.{ext}"
        if p.exists():
            print(f"  Audio exists: {p.name}")
            return p
    url = f"https://www.youtube.com/watch?v={video_id}"
    try:
        with yt_dlp.YoutubeDL({
            "format": "bestaudio[ext=m4a]/bestaudio/best",
            "outtmpl": str(TMP_DIR / f"{label}.%(ext)s"),
            "quiet": False, "noplaylist": True,
        }) as ydl:
            ydl.download([url])
        for ext in ["m4a", "webm", "opus", "mp4"]:
            p = TMP_DIR / f"{label}.{ext}"
            if p.exists():
                return p
    except Exception as e:
        print(f"  Download failed: {e}")
    return None


def cleanup_audio(label: str):
    for ext in ["m4a", "webm", "opus", "mp4", "wav"]:
        p = TMP_DIR / f"{label}.{ext}"
        if p.exists():
            p.unlink()
            print(f"  Deleted: {p.name}")


# ── Gemini 轉錄 ───────────────────────────────────────────────
AUDIO_MIME = {".m4a": "audio/mp4", ".mp4": "audio/mp4",
              ".webm": "audio/webm", ".opus": "audio/ogg"}


def transcribe_with_gemini(client, audio_path: Path, ppt_text: str, title: str) -> str:
    from google.genai import types
    sz_mb = audio_path.stat().st_size / 1024 / 1024
    print(f"  Uploading to Gemini ({sz_mb:.1f} MB)...")

    with tempfile.NamedTemporaryFile(suffix=audio_path.suffix, delete=False) as tmp:
        tmp_path = Path(tmp.name)
    shutil.copyfile(audio_path, tmp_path)

    try:
        mime = AUDIO_MIME.get(audio_path.suffix, "audio/mp4")
        uploaded = client.files.upload(
            file=tmp_path,
            config=types.UploadFileConfig(
                display_name=f"qm_{audio_path.stem[:40]}",
                mime_type=mime,
            ),
        )
    finally:
        try: tmp_path.unlink()
        except Exception: pass

    # 等待 Gemini Files 處理完成
    for _ in range(60):
        fi = client.files.get(name=uploaded.name)
        state = fi.state.name if hasattr(fi.state, "name") else str(fi.state)
        if state == "ACTIVE":
            break
        if state == "FAILED":
            raise RuntimeError("Gemini file processing failed")
        print(f"  File state: {state}, waiting...")
        time.sleep(3)

    context_section = ""
    if ppt_text:
        context_section = f"""以下是本集的投影片內容，作為人名、地名、專有名詞的轉錄參考：
---
{ppt_text}
---

"""

    prompt = f"""你正在聆聽一場繁體中文的讀書會演講，主題是《千面上帝》（The Many Faces of God: A History of Religion in World Perspective）。
講者是張辰瑋，這集的講題是：{title}

{context_section}請將這段音訊完整轉錄成繁體中文文字。

轉錄要求：
1. 使用繁體中文（台灣用語）
2. 適當分段：每個段落 3–6 個句子，話題切換時換段，讓全文清晰易讀
3. 正確使用標點符號（，。？！「」《》等）
4. 完整保留講者的原始用詞，不改寫或摘要
5. 人名、地名、書名、專有名詞以繁體中文呈現
6. 引用詩篇、經文、詩句等，完整轉錄
7. 不加時間戳記
8. 直接輸出轉錄內容，不加前言或後記

請開始："""

    print("  Sending to Gemini for transcription...")
    resp = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=[uploaded, prompt],
    )

    try: client.files.delete(name=uploaded.name)
    except Exception: pass

    return resp.text.strip()


# ── Supabase ──────────────────────────────────────────────────
import requests as _req


def upsert_transcript(episode: int, title: str, content: str,
                      video_date: str | None, youtube_id: str | None) -> bool:
    data: dict = {"project_slug": "qiangmian", "episode": episode,
                  "title": title, "content": content}
    if video_date:
        data["video_date"] = video_date
    if youtube_id:
        data["youtube_id"] = youtube_id
    resp = _req.post(
        f"{SUPABASE_URL}/rest/v1/video_transcripts?on_conflict=project_slug,episode",
        json=data,
        headers={
            "apikey": SERVICE_ROLE_KEY,
            "Authorization": f"Bearer {SERVICE_ROLE_KEY}",
            "Content-Type": "application/json",
            "Prefer": "resolution=merge-duplicates",
        },
        timeout=30,
    )
    return resp.status_code in (200, 201)


# ── 處理單集 ──────────────────────────────────────────────────
def process_episode(client, ep_num: int, entry: dict, ppt_files: list[Path]) -> bool:
    video_id   = entry.get("id", "")
    raw_title  = entry.get("title", f"第{ep_num}集")
    upload_date = entry.get("upload_date", "")
    label      = f"{ep_num:02d} {safe_name(raw_title)}"

    print(f"\n[ep {ep_num}] {raw_title}")
    print(f"  video: {video_id}")

    # PPT 對應：第 ep_num 個 PPT 檔案（按日期排序）
    ppt_text = ""
    if ep_num <= len(ppt_files):
        ppt_path = ppt_files[ep_num - 1]
        print(f"  PPT: {ppt_path.name}")
        try:
            ppt_text = load_ppt_text(ppt_path)
            print(f"  PPT: {len(ppt_text)} chars")
        except Exception as e:
            print(f"  PPT read failed: {e}")

    audio_path = download_audio(video_id, label)
    if not audio_path:
        print("  Audio download failed, skip")
        return False

    try:
        transcript = transcribe_with_gemini(client, audio_path, ppt_text, raw_title)
        print(f"  Transcript: {len(transcript):,} chars")
    except Exception as e:
        print(f"  Transcription error: {e}")
        return False

    if upload_date and len(upload_date) == 8:
        date_str: str | None = f"{upload_date[:4]}-{upload_date[4:6]}-{upload_date[6:]}"
    else:
        date_str = None

    full_content = f"{raw_title}\nEpisode: {ep_num}\n"
    if date_str:
        full_content += f"Date: {date_str}\n"
    full_content += f"\n{transcript}\n"

    print(f"  Upserting to Supabase...")
    ok = upsert_transcript(ep_num, raw_title, full_content, date_str, video_id)
    print(f"  {'✓ OK' if ok else '❌ FAILED'}")

    if ok:
        cleanup_audio(label)

    return ok


# ── main ──────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="宗教史讀書會 Gemini 轉錄")
    g = parser.add_mutually_exclusive_group(required=True)
    g.add_argument("--episode", type=str, metavar="N or N-M",
                   help="集數，例如 1 或 1-5")
    g.add_argument("--all", action="store_true", help="全部集數")
    g.add_argument("--list", action="store_true", help="列出播放清單")
    args = parser.parse_args()

    if not GEMINI_KEYS:
        print("❌ No Gemini API key found in .env"); sys.exit(1)

    # Lazy imports
    try: import yt_dlp
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "yt-dlp", "-q"])
    try: from google import genai
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "google-genai", "-q"])
        from google import genai
    from google import genai

    ppt_files = sorted_ppt_files()

    print("Fetching playlist...")
    entries = fetch_playlist()
    print(f"  {len(entries)} videos in playlist")
    print(f"  {len(ppt_files)} PPT files in Drive")

    if args.list:
        for i, e in enumerate(entries, 1):
            ppt = ppt_files[i - 1].name if i <= len(ppt_files) else "(no PPT)"
            print(f"  {i:2d}. {e.get('title', '?')}")
            print(f"       PPT: {ppt}")
        return

    if args.all:
        targets = list(range(1, len(entries) + 1))
    else:
        ep_str = args.episode
        if "-" in ep_str:
            a, b = ep_str.split("-", 1)
            targets = list(range(int(a), int(b) + 1))
        else:
            targets = [int(ep_str)]

    client = genai.Client(api_key=GEMINI_KEYS[0])

    ok_count = 0
    for i, ep_num in enumerate(targets):
        if ep_num < 1 or ep_num > len(entries):
            print(f"Episode {ep_num} out of range (1–{len(entries)}), skip")
            continue
        success = process_episode(client, ep_num, entries[ep_num - 1], ppt_files)
        if success:
            ok_count += 1
        if i < len(targets) - 1:
            time.sleep(5)

    print(f"\nDone. {ok_count}/{len(targets)} OK")


if __name__ == "__main__":
    main()
