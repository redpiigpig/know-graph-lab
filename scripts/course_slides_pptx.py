# -*- coding: utf-8 -*-
"""課堂簡報（PPTX）。

講義全文在線上與紙本，簡報只放課堂上要投影的骨架：命題、對照、圖表、
討論題。因此投影片內容是「編輯過的重點」而非講義段落的搬運。

每一次上課一份，內容以 DECKS[<次數>] 的資料結構描述，版面由本檔統一渲染。
成品依 docs/repo-hygiene.md 不進 git，輸出到 Drive：
  G:\\我的雲端硬碟\\資料\\知識圖工作室\\教學\\{課程資料夾}\\簡報\\

用法：
  python scripts/course_slides_pptx.py          # 第 1 次
  python scripts/course_slides_pptx.py 1 2
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

DRIVE = Path(r'G:\我的雲端硬碟\資料\知識圖工作室\教學')
# 一門課一個資料夾；`--course=sl` 切到宗教系國文講義。
COURSES = {
    'wr': '115-1_世界宗教文化導論',
    'sl': '宗教系國文講義',
    'ch': '115-2_基督宗教概論',
}
FOLDER = COURSES['wr']
IMGDIR = DRIVE / FOLDER / '簡報' / '圖片'


def load_manifest():
    f = IMGDIR / '_manifest.json'
    return json.loads(f.read_text(encoding='utf-8')) if f.exists() else {}


MANIFEST = load_manifest()
USED = []          # 本份簡報實際用到的圖，供「圖片出處」頁使用

HEI = '微軟正黑體'
KAI = '標楷體'
NAVY = RGBColor(0x1F, 0x39, 0x64)
GOLD = RGBColor(0xB2, 0x8B, 0x3C)
INK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x6B, 0x72, 0x80)
PALE = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Cm(33.87), Cm(19.05)   # 16:9


# ── 版面工具 ────────────────────────────────────────────────────────────────
def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


# 「」內六字以上視為引語（學者的話、經文），排標楷體；
# 六字以下多半是專名或術語標記（「基督教」「已然—未然」），維持原字體。
QUOTE_RE = re.compile(r'[「『][^「」『』]{6,}[」』]')


def _spans(text):
    out, at = [], 0
    for m in QUOTE_RE.finditer(text):
        if m.start() > at:
            out.append((text[at:m.start()], False))
        out.append((m.group(), True))
        at = m.end()
    if at < len(text):
        out.append((text[at:], False))
    return out or [(text, False)]


def put(tf, text, size, font=HEI, bold=False, color=INK, space_after=6,
        align=PP_ALIGN.LEFT, first=False, line=1.25, indent=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if indent:
        p.level = min(indent, 4)
    for chunk, is_quote in _spans(text):
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.name = KAI if (is_quote and font != KAI) else font
        r.font.bold = bold
        r.font.color.rgb = color
    return p


CM_PT = 28.35


def fit(items, width_cm, height_cm, sizes, spaces, line=1.3, indent_cm=(0, 0.9, 1.6)):
    """估算這批條目實際佔幾行，回傳縮放係數（最小 0.72）。

    中日文一個字約等於一個字級的寬度，因此每行字數 ≈ 可用寬度 ÷ 字級。
    只縮小、不放大——版面預設就是給內容少的頁看的。
    """
    # PowerPoint 的中文行高不是「字級×行距」，而是再乘上字型的行距係數
    # （實測約 1.22）。低估的話文字會被切掉或壓到頁尾，所以這裡取 1.62。
    LINE = 1.75 / 1.3 * line
    total = 0.0
    for it in items:
        lvl, txt = (it if isinstance(it, tuple) else (0, it))
        if not txt:
            total += sizes[0] * 0.5
            continue
        avail = (width_cm - indent_cm[min(lvl, 2) if lvl != 3 else 0]) * CM_PT
        per = max(8, int(avail / sizes[lvl]))
        rows = -(-(len(txt) + 2) // per)          # ＋2 是行首的項目符號
        total += rows * sizes[lvl] * LINE + spaces[lvl]
    return min(1.0, max(0.55, height_cm * CM_PT / total)) if total else 1.0


def band(slide, color, x, y, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_title(slide, title, sub=None):
    band(slide, NAVY, Cm(0), Cm(0), W, Cm(0.32))
    tf = textbox(slide, Cm(1.5), Cm(0.85), W - Cm(3.0), Cm(2.2))
    put(tf, title, 34, bold=True, color=NAVY, first=True, space_after=2)
    if sub:
        put(tf, sub, 17, color=GRAY, space_after=0)
    band(slide, GOLD, Cm(1.5), Cm(3.80), Cm(3.0), Cm(0.12))


def footer(slide, n, label):
    tf = textbox(slide, Cm(1.5), H - Cm(1.15), W - Cm(3.0), Cm(0.8))
    put(tf, f'{label}　　{n}', 11, color=GRAY, first=True, space_after=0)


# ── 各種投影片 ──────────────────────────────────────────────────────────────
def s_cover(prs, d):
    s = blank(prs)
    band(s, NAVY, Cm(0), Cm(0), W, H)
    tf = textbox(s, Cm(3.0), Cm(4.2), W - Cm(6.0), Cm(1.0))
    put(tf, d['kicker'], 15, color=RGBColor(0xC9, 0xD3, 0xE4), first=True, space_after=0)
    tf = textbox(s, Cm(3.0), Cm(5.5), W - Cm(6.0), Cm(2.4))
    put(tf, d['title'], 44, font=KAI, bold=True, color=WHITE, first=True, space_after=0)
    band(s, GOLD, Cm(3.0), Cm(8.5), Cm(3.4), Cm(0.12))
    tf2 = textbox(s, Cm(3.0), Cm(9.5), W - Cm(6.0), Cm(6.0))
    put(tf2, d['subtitle'], 20, font=KAI, color=RGBColor(0xE3, 0xC9, 0x8A),
        first=True, space_after=24)
    for line in d['meta']:
        put(tf2, line, 14, color=RGBColor(0xC9, 0xD3, 0xE4), space_after=6)
    return s


def s_section(prs, no, title, lines):
    s = blank(prs)
    band(s, PALE, Cm(0), Cm(0), W, H)
    band(s, NAVY, Cm(0), Cm(0), Cm(0.5), H)
    tf = textbox(s, Cm(3.0), Cm(4.8), W - Cm(6.0), Cm(9.5))
    put(tf, no, 19, bold=True, color=GOLD, first=True, space_after=12)
    put(tf, title, 42, font=KAI, bold=True, color=NAVY, space_after=20)
    for ln in lines:
        put(tf, ln, 18, color=GRAY, space_after=8)
    return s


def s_big(prs, text, sub=None):
    s = blank(prs)
    band(s, NAVY, Cm(0), Cm(0), W, H)
    tf = textbox(s, Cm(3.2), Cm(2.6), W - Cm(6.4), H - Cm(5.2), anchor=MSO_ANCHOR.MIDDLE)
    put(tf, text, 36, font=KAI, bold=True, color=WHITE, first=True,
        align=PP_ALIGN.CENTER, line=1.45, space_after=18)
    if sub:
        put(tf, sub, 18, color=RGBColor(0xE3, 0xC9, 0x8A), align=PP_ALIGN.CENTER)
    return s


def s_bullets(prs, title, bullets, sub=None):
    s = blank(prs)
    slide_title(s, title, sub)
    w, h = 30.9, 13.1
    tf = textbox(s, Cm(1.5), Cm(4.45), Cm(w), Cm(h))
    base = {0: 23.0, 1: 19.0, 2: 16.5, 3: 24.0}
    sp = {0: 11, 1: 7, 2: 5, 3: 12}
    k = fit(bullets, w, h * 0.96, base, sp)
    # 內容明顯偏少（六成高度就裝得下）就垂直置中，不要下半頁整片空白
    if k >= 1.0 and fit(bullets, w, h * 0.80, base, sp) >= 1.0:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    firstdone = False
    for b in bullets:
        lvl, txt = (b if isinstance(b, tuple) else (0, b))
        if txt == '':
            put(tf, ' ', 9 * k, first=not firstdone, space_after=0)
            firstdone = True
            continue
        color = {0: INK, 1: RGBColor(0x3A, 0x3A, 0x3A), 2: GRAY, 3: INK}[lvl]
        mark = {0: '▍', 1: '‧', 2: '－', 3: ''}[lvl]
        put(tf, (f'{mark} ' if mark else '') + txt, base[lvl] * k, color=color,
            bold=(lvl == 0), first=not firstdone, space_after=sp[lvl] * k,
            indent=min(lvl, 2) if lvl != 3 else 0, line=1.3)
        firstdone = True
    return s


def s_two(prs, title, left, right, sub=None):
    s = blank(prs)
    slide_title(s, title, sub)
    colw = (W - Cm(3.8)) / 2
    for i, (head, items) in enumerate((left, right)):
        x = Cm(1.5) + i * (colw + Cm(0.8))
        band(s, NAVY if i == 0 else GOLD, x, Cm(4.5), colw, Cm(1.0))
        tfh = textbox(s, x + Cm(0.35), Cm(4.68), colw - Cm(0.7), Cm(0.8))
        put(tfh, head, 18, bold=True, color=WHITE, first=True, space_after=0)
        cw = colw / 360000 / 10 - 0.7
        base = {0: 18.5, 1: 15.5, 2: 15.5}
        k = max(0.68, fit(items, cw, 13.0, base, {0: 8, 1: 6, 2: 6}))
        tf = textbox(s, x + Cm(0.35), Cm(5.8), colw - Cm(0.7), H - Cm(7.15))
        for j, it in enumerate(items):
            lvl, txt = (it if isinstance(it, tuple) else (0, it))
            put(tf, ('‧ ' if lvl == 0 else '　－ ') + txt,
                base[min(lvl, 1)] * k,
                color=INK if lvl == 0 else GRAY,
                first=(j == 0), space_after=8 * k, line=1.3)
    return s


def s_table(prs, title, headers, rows, sub=None, note=None, widths=None):
    s = blank(prs)
    slide_title(s, title, sub)
    top = Cm(4.5)
    tbl = s.shapes.add_table(len(rows) + 1, len(headers),
                             Cm(1.5), top, W - Cm(3.0), Cm(1.0)).table
    # 列高給最小值，讓 PowerPoint 依內容自動撐開（不要平均攤滿整頁）
    for i, r in enumerate(tbl.rows):
        r.height = Cm(1.05) if i == 0 else Cm(0.85)
    if widths:
        total = sum(widths)
        for i, w in enumerate(widths):
            tbl.columns[i].width = int((W - Cm(3.0)) * w / total)
    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = ''
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame; tf.word_wrap = True
        put(tf, h, 15, bold=True, color=WHITE, first=True, space_after=0,
            align=PP_ALIGN.CENTER)
    tsize = 15 if len(rows) <= 5 else (13.5 if len(rows) <= 7 else 12)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = ''
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if ri % 2 == 0 else PALE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Cm(0.12)
            tf = c.text_frame; tf.word_wrap = True
            put(tf, str(val), tsize, color=INK, first=True, space_after=0,
                line=1.15)
    if note:
        tf = textbox(s, Cm(1.5), H - Cm(2.6), W - Cm(3.0), Cm(1.3))
        put(tf, note, 12.5, color=GRAY, first=True, space_after=0)
    return s


def place_image(slide, key, x, y, w, h):
    """把圖等比塞進 (x,y,w,h) 的框裡並置中；圖不存在就回傳 False。"""
    m = MANIFEST.get(key)
    if not m:
        return False
    f = IMGDIR / m['file']
    if not f.exists():
        return False
    pic = slide.shapes.add_picture(str(f), x, y, width=w)
    if pic.height > h:                      # 太高就改用高度縮，維持比例
        ratio = pic.width / pic.height
        pic.height = h
        pic.width = int(h * ratio)
    pic.left = x + int((w - pic.width) / 2)
    pic.top = y + int((h - pic.height) / 2)
    if key not in USED:
        USED.append(key)
    return True


def caption(slide, text, x, y, w):
    tf = textbox(slide, x, y, w, Cm(1.3))
    put(tf, text, 13, color=GRAY, first=True, space_after=0, align=PP_ALIGN.CENTER)


def s_photo(prs, title, key, cap=None, sub=None):
    s = blank(prs)
    slide_title(s, title, sub)
    box_y, box_h = Cm(4.45), H - Cm(6.75)
    if not place_image(s, key, Cm(1.5), box_y, W - Cm(3.0), box_h):
        put(textbox(s, Cm(1.5), box_y, W - Cm(3.0), Cm(2)), f'（缺圖：{key}）',
            16, color=GRAY, first=True)
    if cap:
        caption(s, cap, Cm(1.5), H - Cm(2.1), W - Cm(3.0))
    return s


def s_gallery(prs, title, items, sub=None):
    """items: [(圖 key, 說明), ...]，二至四張並排。"""
    s = blank(prs)
    slide_title(s, title, sub)
    n = len(items)
    gap = Cm(0.6)
    colw = int((W - Cm(3.0) - gap * (n - 1)) / n)
    top, boxh = Cm(4.5), H - Cm(7.65)
    for i, (key, label) in enumerate(items):
        x = Cm(1.5) + i * (colw + gap)
        if not place_image(s, key, x, top, colw, boxh):
            put(textbox(s, x, top, colw, Cm(2)), f'（缺圖：{key}）', 13,
                color=GRAY, first=True)
        tf = textbox(s, x, top + boxh + Cm(0.25), colw, Cm(1.9))
        put(tf, label, 13.5, color=INK, first=True, space_after=0,
            align=PP_ALIGN.CENTER, line=1.25)
    return s


def s_imgbullets(prs, title, bullets, key, sub=None, cap=None):
    """左文右圖。"""
    s = blank(prs)
    slide_title(s, title, sub)
    tw = 17.0
    textw = Cm(tw)
    imgx = Cm(1.5) + textw + Cm(0.7)
    imgw = W - Cm(1.5) - imgx
    h = 13.1
    tf = textbox(s, Cm(1.5), Cm(4.45), textw, Cm(h))
    base = {0: 20.0, 1: 17.0, 2: 14.5, 3: 21.0}
    sp = {0: 10, 1: 7, 2: 5, 3: 11}
    k = fit(bullets, tw, h * 0.96, base, sp)
    firstdone = False
    for b in bullets:
        lvl, txt = (b if isinstance(b, tuple) else (0, b))
        if txt == '':
            put(tf, ' ', 9 * k, first=not firstdone, space_after=0)
            firstdone = True
            continue
        color = {0: INK, 1: RGBColor(0x3A, 0x3A, 0x3A), 2: GRAY, 3: INK}[lvl]
        mark = {0: '▍', 1: '‧', 2: '－', 3: ''}[lvl]
        put(tf, (f'{mark} ' if mark else '') + txt, base[lvl] * k, color=color,
            bold=(lvl == 0), first=not firstdone, space_after=sp[lvl] * k,
            indent=min(lvl, 2) if lvl != 3 else 0, line=1.3)
        firstdone = True
    boxh = H - Cm(6.4) if cap else H - Cm(5.4)
    place_image(s, key, imgx, Cm(4.45), imgw, boxh)
    if cap:
        caption(s, cap, imgx, H - Cm(2.1), imgw)
    return s


def s_credits(prs, label):
    """圖片出處頁：課堂投影用圖必須標明作者與授權。

    條目多就分頁——擠成一頁小字會看不清，而這一頁是授權聲明，
    看不清等於沒標。
    """
    if not USED:
        return []
    per, out = 11, []
    for start in range(0, len(USED), per):
        s = blank(prs)
        slide_title(s, '圖片出處',
                    '本份簡報用圖均取自維基共享資源，授權為公有領域或 CC')
        tf = textbox(s, Cm(1.5), Cm(4.45), W - Cm(3.0), H - Cm(6.15))
        for i, k in enumerate(USED[start:start + per]):
            m = MANIFEST.get(k, {})
            name = m.get('title', k)[5:]            # 去掉 'File:'
            line = f'{name}　—　{m.get("license", "")}'
            if m.get('author'):
                line += f'　／　{m["author"][:44]}'
            put(tf, line, 12.5, color=GRAY, first=(i == 0),
                space_after=6, line=1.25)
        out.append(s)
    return out


RENDER = {'cover': s_cover, 'section': s_section, 'big': s_big,
          'bullets': s_bullets, 'two': s_two, 'table': s_table,
          'photo': s_photo, 'gallery': s_gallery, 'imgbullets': s_imgbullets}



# ── 併頁：不要「一頁只有一句話」 ────────────────────────────────────────
BULLETY = ('bullets', 'imgbullets')


def _big_lines(b):
    # level 3＝導言／引文行：不掛項目符號，出處自成一行
    lines = [(3, ln) for ln in b[1].split(chr(10)) if ln.strip()]
    if len(b) > 2 and b[2]:
        lines.append((1, b[2]))
    return lines


def _swap_bullets(it, newb):
    return (it[0], it[1], newb) + tuple(it[3:])


def _to_bullets(b):
    """真的無處可併時，讓它自己成為一張有標題有內容的條列頁。"""
    lines = [ln for ln in b[1].split('\n') if ln.strip()]
    title = lines[0].rstrip('。，、')
    rest = [(0, ln) for ln in lines[1:]]
    if len(b) > 2 and b[2]:
        rest.append((1, b[2]))
    return ('bullets', title, rest or [(1, '')])


def fold_bigs(slides):
    out, pending = [], []
    for it in slides:
        if it[0] == 'big':
            pending.append(it)
            continue
        if pending:
            extra = [x for b in pending for x in _big_lines(b)]
            if it[0] in BULLETY:
                it = _swap_bullets(it, extra + [''] + list(it[2]))
            elif it[0] == 'section':
                lines = list(it[3]) if len(it) > 3 else []
                for b in pending:
                    lines += [ln for ln in b[1].split('\n') if ln.strip()]
                    if len(b) > 2 and b[2]:
                        lines.append(b[2])
                it = (it[0], it[1], it[2], lines)
            elif out and out[-1][0] in BULLETY:
                out[-1] = _swap_bullets(out[-1], list(out[-1][2]) + [''] + extra)
            else:
                out.extend(_to_bullets(b) for b in pending)
            pending = []
        out.append(it)
    for b in pending:
        if out and out[-1][0] in BULLETY:
            out[-1] = _swap_bullets(out[-1], list(out[-1][2]) + [''] + _big_lines(b))
        else:
            out.append(_to_bullets(b))
    return out


def split_long(slides):
    """一頁塞不下就拆成兩頁，不要把字縮到看不清。"""
    base = {0: 23.0, 1: 19.0, 2: 16.5, 3: 24.0}
    sp = {0: 11, 1: 7, 2: 5, 3: 12}
    ibase = {0: 20.0, 1: 17.0, 2: 14.5, 3: 21.0}
    isp = {0: 10, 1: 7, 2: 5, 3: 11}
    out = []
    for it in slides:
        if it[0] == 'bullets':
            k = fit(it[2], 30.9, 13.1 * 0.96, base, sp)
        elif it[0] == 'imgbullets':
            k = fit(it[2], 17.0, 13.1 * 0.96, ibase, isp)
        else:
            out.append(it)
            continue
        if k >= 0.78:
            out.append(it)
            continue
        items = list(it[2])
        # 從中間往後找第一個第一層項目當切點，避免把子項目跟標題拆開
        half = len(items) // 2
        cut = next((i for i in range(half, len(items))
                    if not isinstance(items[i], tuple) or items[i][0] == 0), half)
        if it[0] == 'bullets':
            rest = tuple(it[3:])
            out.append(('bullets', it[1], items[:cut]) + rest)
            out.append(('bullets', it[1] + '（續）', items[cut:]) + rest)
        else:
            # 圖留在第一頁，續頁走純文字（整頁寬，字才放得大）
            out.append((it[0], it[1], items[:cut]) + tuple(it[3:]))
            out.append(('bullets', it[1] + '（續）', items[cut:]))
    return out

def build(deck):
    global USED
    USED = []
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for i, item in enumerate(split_long(fold_bigs(deck['slides']))):
        kind, args = item[0], list(item[1:])
        kw = args.pop() if args and isinstance(args[-1], dict) and kind != 'cover' else {}
        s = RENDER[kind](prs, *args, **kw)
        if kind not in ('cover', 'section', 'big'):
            footer(s, i, deck['footer'])
    for c in s_credits(prs, deck['footer']):
        footer(c, len(prs.slides._sldIdLst) - 1, deck['footer'])
    outdir = DRIVE / FOLDER / '簡報'
    outdir.mkdir(parents=True, exist_ok=True)
    out = outdir / deck['filename']
    prs.save(out)
    return out, len(prs.slides._sldIdLst)


if __name__ == '__main__':
    sys.stdout.reconfigure(encoding='utf-8')
    args = sys.argv[1:]
    course = next((a.split('=')[1] for a in args if a.startswith('--course=')), 'wr')
    if course != 'wr':
        FOLDER = COURSES[course]
        IMGDIR = DRIVE / FOLDER / '簡報' / '圖片'
        MANIFEST = load_manifest()
    if course in ('sl', 'ch'):
        if course == 'sl':
            from course_slides_data_sl import DECKS_SL as D
        else:
            from course_slides_data_ch import DECKS_CH as D
        DECKS = D
        for n in [a for a in args if not a.startswith('--')] or sorted(DECKS):
            out, cnt = build(DECKS[int(n)])
            print(f'✔ {out}　（{cnt} 張）')
        raise SystemExit
    from course_slides_data import DECKS
    from course_slides_data2 import DECKS2
    from course_slides_data3 import DECKS3
    from course_slides_data4 import DECKS4
    DECKS = {**DECKS, **DECKS2, **DECKS3, **DECKS4}
    for n in ([a for a in args if not a.startswith('--')] or sorted(DECKS)):
        out, cnt = build(DECKS[int(n)])
        print(f'✔ {out}　（{cnt} 張）')
