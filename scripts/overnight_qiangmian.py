#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
千面上帝 整晚自動化：用 local Whisper 補轉錄缺失/截斷的集數，用 Haiku（OAuth）
潤稿、加 PPT 參考書目，全部寫回 Supabase。

Phase 1（快）：對 DB 內已完整的 Gemini 轉錄做 Haiku 潤稿
Phase 2（慢）：對缺失或截斷的集數，先 Whisper 轉錄再 Haiku 潤稿

執行：python -u scripts/overnight_qiangmian.py
"""
import json
import os
import re
import subprocess
import sys
import time
from pathlib import Path

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

REPO_DIR = Path(__file__).parent.parent
PPT_DIR  = Path("G:/我的雲端硬碟/創作/千面上帝/宗教史讀書會")
TMP_DIR  = REPO_DIR / "_tmp_audio" / "qiangmian"
TMP_DIR.mkdir(parents=True, exist_ok=True)

PLAYLIST_URL = "https://www.youtube.com/playlist?list=PLNdU3g_-OSshfnyOakO5exMMvnSNeuIjZ"

# ── .env ────────────────────────────────────────────────────
_env: dict[str, str] = {}
for _line in (REPO_DIR / ".env").read_text(encoding="utf-8").splitlines():
    if "=" in _line and not _line.startswith("#"):
        k, _, v = _line.partition("=")
        _env[k.strip()] = v.strip().strip('"')

SUPABASE_URL = _env["SUPABASE_URL"]
SERVICE_KEY  = _env["SUPABASE_SERVICE_ROLE_KEY"]

# ── 哪些集要做什麼 ──────────────────────────────────────────
# POLISH_ONLY: DB 已有完整 Gemini 轉錄，只跑 Haiku 潤稿
POLISH_ONLY_DEFAULT = [3, 4, 5, 6, 8, 9, 10, 11]
# TRANSCRIBE: 需要 Whisper 重新轉錄（loop / truncated / missing / 503 失敗）
TRANSCRIBE_DEFAULT  = [2, 7, 12, 13, 14, 15, 16, 17, 18, 19, 20, 21, 22, 23, 24, 25]

# YouTube 影片 → PPT 章節對應（按章節內容人工比對）
# 不能用 ppt_files[ep-1] 機械對應，因為某些章節有 PPT 但沒影片，
# 或者一份 PPT 講三集（如第十二章上）。
# Key: episode 編號（YouTube playlist 順序）；Value: PPT index（依日期排序，1-based）
EP_TO_PPT_IDX = {
    1: 1, 2: 2, 3: 3, 4: 4,
    5: 6,   # ep 5 = 第六章（跳過 PPT 5 = 第五章，因為沒影片）
    6: 7, 7: 8, 8: 9, 9: 10, 10: 11, 11: 12, 12: 13,
    13: 14, 14: 14, 15: 14,  # 第十二章上：三集共用一份 PPT
    16: 15,
    17: 16, 18: 17, 19: 18, 20: 19,
    21: 21,  # 第十五章下（跳過 PPT 20 = 第十五章上，因為沒影片）
    22: 22,
    23: 26,  # 第十七章中 ≈ 第十七章下（PPT 沒有「中」這個版本，best guess）
    24: 28, 25: 29,
}

# 命令列覆寫（--polish 1,2,3 --transcribe 4-10）
def _parse_eps(s: str) -> list[int]:
    out: list[int] = []
    for part in s.split(","):
        part = part.strip()
        if not part: continue
        if "-" in part:
            a, b = part.split("-", 1)
            out.extend(range(int(a), int(b) + 1))
        else:
            out.append(int(part))
    return out

POLISH_ONLY = POLISH_ONLY_DEFAULT
TRANSCRIBE  = TRANSCRIBE_DEFAULT
_argv = sys.argv[1:]
i = 0
while i < len(_argv):
    if _argv[i] == "--polish" and i + 1 < len(_argv):
        POLISH_ONLY = _parse_eps(_argv[i+1]); i += 2
    elif _argv[i] == "--transcribe" and i + 1 < len(_argv):
        TRANSCRIBE = _parse_eps(_argv[i+1]); i += 2
    elif _argv[i] == "--no-polish":
        POLISH_ONLY = []; i += 1
    elif _argv[i] == "--no-transcribe":
        TRANSCRIBE = []; i += 1
    else:
        i += 1

HAIKU_MODEL = "claude-haiku-4-5-20251001"
WHISPER_SIZE = "small"
CHUNK_SECONDS = 1800  # 30 分鐘


# ── Anthropic via OAuth（與 ocr_with_claude_cli.py 共用 token）──
def make_anthropic_client():
    import anthropic
    api_key = _env.get("ANTHROPIC_API_KEY") or os.environ.get("ANTHROPIC_API_KEY")
    if api_key:
        return anthropic.Anthropic(api_key=api_key)
    cred_path = Path(os.environ.get("USERPROFILE", os.environ.get("HOME", ""))) / ".claude" / ".credentials.json"
    if cred_path.exists():
        with cred_path.open(encoding="utf-8") as f:
            creds = json.load(f)
        token = creds.get("claudeAiOauth", {}).get("accessToken", "")
        if token:
            return anthropic.Anthropic(auth_token=token)
    raise RuntimeError("No Anthropic credentials")


# ── PPT ─────────────────────────────────────────────────────
def sorted_ppt_files() -> list[Path]:
    files = list(PPT_DIR.glob("*.pptx"))
    files.sort(key=lambda p: p.name[:10])
    return files


def load_ppt_text(ppt_path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(ppt_path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            out.append(f"[投影片{i}] " + " / ".join(parts))
    return "\n".join(out)


def load_ppt_bibliography(ppt_path: Path) -> str:
    """Read last slide(s) which usually contain the 參考書目 list."""
    from pptx import Presentation
    prs = Presentation(str(ppt_path))
    slides = list(prs.slides)
    # Walk from last slide backward, find one containing 參考書目 marker
    for slide in reversed(slides[-3:]):
        text = "\n".join(s.text.strip() for s in slide.shapes
                         if hasattr(s, "text") and s.text.strip())
        if "參考書目" in text or "參考資料" in text or "書目" in text:
            return text
    # Fallback: take last slide
    last = slides[-1]
    return "\n".join(s.text.strip() for s in last.shapes
                     if hasattr(s, "text") and s.text.strip())


# ── Supabase ────────────────────────────────────────────────
import requests as _req

def fetch_row(episode: int) -> dict | None:
    r = _req.get(
        f"{SUPABASE_URL}/rest/v1/video_transcripts",
        params={"project_slug": "eq.qiangmian", "episode": f"eq.{episode}",
                "select": "id,title,content,video_date,youtube_id,ppt_r2_key"},
        headers={"apikey": SERVICE_KEY, "Authorization": f"Bearer {SERVICE_KEY}"},
        timeout=30,
    )
    rows = r.json() if r.status_code == 200 else []
    return rows[0] if rows else None


def upsert_row(episode: int, title: str, content: str,
               video_date: str | None, youtube_id: str | None,
               ppt_r2_key: str | None) -> bool:
    data = {"project_slug": "qiangmian", "episode": episode,
            "title": title, "content": content}
    if video_date:  data["video_date"] = video_date
    if youtube_id:  data["youtube_id"] = youtube_id
    if ppt_r2_key:  data["ppt_r2_key"] = ppt_r2_key
    r = _req.post(
        f"{SUPABASE_URL}/rest/v1/video_transcripts?on_conflict=project_slug,episode",
        json=data,
        headers={"apikey": SERVICE_KEY, "Authorization": f"Bearer {SERVICE_KEY}",
                 "Content-Type": "application/json",
                 "Prefer": "resolution=merge-duplicates"},
        timeout=30,
    )
    return r.status_code in (200, 201)


# ── Playlist ────────────────────────────────────────────────
def fetch_playlist() -> list[dict]:
    import yt_dlp
    with yt_dlp.YoutubeDL({"quiet": True, "extract_flat": True,
                            "cookiesfrombrowser": ("firefox",)}) as ydl:
        info = ydl.extract_info(PLAYLIST_URL, download=False)
    return info.get("entries", [])


def safe_name(s: str) -> str:
    return re.sub(r'[\\/:*?"<>|]', '_', s)[:80].strip()


# ── 下載音訊（已存在則重用）────────────────────────────────
def download_audio(video_id: str, label: str, max_attempts: int = 3) -> Path | None:
    import yt_dlp
    for ext in ["m4a", "webm", "opus", "mp4"]:
        p = TMP_DIR / f"{label}.{ext}"
        if p.exists():
            print(f"  Audio cached: {p.name}")
            return p
    url = f"https://www.youtube.com/watch?v={video_id}"
    for attempt in range(1, max_attempts + 1):
        try:
            with yt_dlp.YoutubeDL({
                "format": "bestaudio[ext=m4a]/bestaudio/best",
                "outtmpl": str(TMP_DIR / f"{label}.%(ext)s"),
                "quiet": False, "noplaylist": True,
                "socket_timeout": 30, "retries": 3,
                "cookiesfrombrowser": ("firefox",),
            }) as ydl:
                ydl.download([url])
            for ext in ["m4a", "webm", "opus", "mp4"]:
                p = TMP_DIR / f"{label}.{ext}"
                if p.exists():
                    return p
        except Exception as e:
            print(f"  Download attempt {attempt} failed: {e}")
            for ext in ["m4a", "webm", "opus", "mp4", "part"]:
                p = TMP_DIR / f"{label}.{ext}"
                if p.exists():
                    p.unlink()
            if attempt < max_attempts:
                time.sleep(15 * attempt)
    return None


# ── Whisper（local）──────────────────────────────────────────
_whisper_model = None

def get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        from faster_whisper import WhisperModel
        print(f"Loading Whisper model ({WHISPER_SIZE}, CPU, int8)...")
        _whisper_model = WhisperModel(WHISPER_SIZE, device="cpu", compute_type="int8")
    return _whisper_model


def to_wav(audio: Path) -> Path:
    wav = audio.with_suffix(".wav")
    if wav.exists():
        return wav
    print("  Converting to WAV...")
    subprocess.call(
        ["ffmpeg", "-i", str(audio), "-ar", "16000", "-ac", "1",
         "-c:a", "pcm_s16le", str(wav), "-y"],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return wav


def get_duration(path: Path) -> float:
    r = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration",
         "-of", "default=noprint_wrappers=1:nokey=1", str(path)],
        capture_output=True, text=True)
    try: return float(r.stdout.strip())
    except: return 0.0


def make_chunk_wav(src: Path, start: int, duration: int, out_path: Path) -> bool:
    if out_path.exists(): return True
    ret = subprocess.call(
        ["ffmpeg", "-i", str(src), "-ss", str(start), "-t", str(duration),
         "-ar", "16000", "-ac", "1", "-c:a", "pcm_s16le", str(out_path), "-y"],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return ret == 0


def transcribe_audio_whisper(audio_path: Path) -> str:
    """Returns raw transcript text (no segments structure)."""
    model = get_whisper_model()
    duration = get_duration(audio_path)
    print(f"  Audio duration: {duration/60:.1f} min")

    all_text: list[str] = []
    if duration <= CHUNK_SECONDS:
        wav = to_wav(audio_path)
        print("  Whisper transcribing...")
        segs, info = model.transcribe(str(wav), language="zh", beam_size=1,
                                       vad_filter=True,
                                       vad_parameters={"min_silence_duration_ms": 1000})
        all_text = [s.text for s in segs]
        wav.unlink(missing_ok=True)
    else:
        n = int(duration // CHUNK_SECONDS) + 1
        print(f"  Long audio, splitting into {n} chunks of 30 min")
        for c in range(n):
            start = c * CHUNK_SECONDS
            chunk_wav = TMP_DIR / f"{audio_path.stem}_chunk{c:02d}.wav"
            print(f"  [{c+1}/{n}] {start//60}-{(start+CHUNK_SECONDS)//60} min...")
            if not make_chunk_wav(audio_path, start, CHUNK_SECONDS, chunk_wav):
                print("    chunk ffmpeg failed"); continue
            try:
                segs, info = model.transcribe(str(chunk_wav), language="zh", beam_size=1,
                                               vad_filter=True,
                                               vad_parameters={"min_silence_duration_ms": 1000})
                chunk_text = "".join(s.text for s in segs)
                all_text.append(chunk_text)
            except Exception as e:
                print(f"    transcribe error: {e}")
            chunk_wav.unlink(missing_ok=True)

    return "\n\n".join(t.strip() for t in all_text if t.strip())


# ── Haiku 潤稿 ──────────────────────────────────────────────
POLISH_PROMPT = """你正在處理一份《千面上帝》宗教史讀書會的逐字稿。講者是張辰瑋，本集的主題：{title}

以下是本集投影片內容（人名、地名、書名、專有名詞的正確寫法參考）：
---
{ppt_text}
---

以下是本集的原始逐字稿（可能來自 Gemini 或 Whisper）：
---
{raw}
---

請完成以下整理：
1. 統一為「繁體中文（台灣用語）」
2. 刪除口語填充詞：就是說、這樣子、嘛、那個（語助詞）、口吃重複（「我、我」）、場務對話（「請問看得到投影片嗎？」之類）
3. 合併重複句子，但保留講者所有實質內容（不要摘要、不要改寫成自己的話）
4. 適當分段：每段 3–6 句，主題切換時換段
5. 加入 ## 小標：全文 8–12 個小標，每個 2–4 字（例：## 起源、## 神話、## 薩滿）
6. 標點：逗號、句號、問號、引號（「」）、書名號（《》）正確使用
7. 人名、地名、書名、專有名詞優先採用 PPT 中的寫法

只輸出潤稿後的正文（## 小標 + 段落），不要任何前言、說明、結語。"""


CHUNK_THRESHOLD = 20000  # raw chars; above this we split for polish

def _polish_one(client, raw: str, title: str, ppt_text: str, position: str) -> str:
    """position: 'full', 'first half', 'second half'"""
    extra = ""
    if position != "full":
        extra = f"\n\n注意：這是逐字稿的「{position}」部分，請只整理這段，不要加開頭/結尾的綜述。"
    prompt = POLISH_PROMPT.format(title=title, ppt_text=ppt_text[:4000], raw=raw) + extra
    for attempt in range(1, 5):
        try:
            chunks: list[str] = []
            with client.messages.stream(
                model=HAIKU_MODEL,
                max_tokens=32000,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    chunks.append(text)
            return "".join(chunks).strip()
        except Exception as e:
            err = str(e)
            print(f"  Haiku attempt {attempt} error: {err[:200]}")
            if attempt < 4:
                time.sleep(10 * attempt)
            else:
                raise


def polish_with_haiku(client, raw_transcript: str, title: str, ppt_text: str) -> str:
    if len(raw_transcript) <= CHUNK_THRESHOLD:
        return _polish_one(client, raw_transcript, title, ppt_text, "full")

    # Split at the closest paragraph break to the midpoint
    mid = len(raw_transcript) // 2
    paragraphs = raw_transcript.split("\n\n")
    cum = 0
    split_idx = 0
    for i, p in enumerate(paragraphs):
        cum += len(p) + 2
        if cum >= mid:
            split_idx = i + 1
            break
    first = "\n\n".join(paragraphs[:split_idx]).strip()
    second = "\n\n".join(paragraphs[split_idx:]).strip()
    print(f"  Chunked polish: {len(first):,} + {len(second):,} chars")

    p1 = _polish_one(client, first, title, ppt_text, "前半段")
    print(f"  first half polished: {len(p1):,} chars")
    p2 = _polish_one(client, second, title, ppt_text, "後半段")
    print(f"  second half polished: {len(p2):,} chars")
    return p1.rstrip() + "\n\n" + p2.lstrip()


# ── 組合最終 content ────────────────────────────────────────
def build_final_content(title: str, episode: int, video_date: str | None,
                        polished: str, bibliography_block: str) -> str:
    out = [title, f"Episode: {episode}"]
    if video_date:
        out.append(f"Date: {video_date}")
    out += ["", polished.strip(), "", "---", "", "*參考書目*", "", bibliography_block.strip()]
    return "\n".join(out)


def extract_polished_raw(content: str) -> str:
    """Strip the title/episode/date header and trailing bibliography from existing content."""
    lines = content.splitlines()
    # Skip leading title + Episode/Date lines
    i = 0
    while i < len(lines) and (i < 4 or lines[i].strip() == ""):
        if lines[i].startswith("Episode:") or lines[i].startswith("Date:"):
            i += 1; continue
        if i == 0:
            i += 1; continue
        if lines[i].strip() == "":
            i += 1
            break
        i += 1
    body = "\n".join(lines[i:])
    # Remove trailing bibliography (--- onwards)
    if "\n---\n" in body:
        body = body.split("\n---\n", 1)[0]
    return body.strip()


# ── 主流程 ──────────────────────────────────────────────────
def cleanup_audio(label: str):
    for ext in ["m4a", "webm", "opus", "mp4", "wav"]:
        p = TMP_DIR / f"{label}.{ext}"
        if p.exists():
            try:
                p.unlink()
                print(f"  deleted: {p.name}")
            except Exception:
                pass


def process_episode(client, ep: int, entries: list, ppt_files: list[Path],
                    do_transcribe: bool) -> bool:
    if ep < 1 or ep > len(entries):
        print(f"[ep {ep}] out of range, skip")
        return False

    entry = entries[ep - 1]
    video_id = entry.get("id", "")
    raw_title = entry.get("title", f"第{ep}集")
    label = f"{ep:02d} {safe_name(raw_title)}"

    print(f"\n[ep {ep}] {raw_title}  (transcribe={do_transcribe})")

    # PPT 對應（用章節對應表，不可用 ppt_files[ep-1]）
    ppt_idx = EP_TO_PPT_IDX.get(ep)
    ppt_text = ""
    bibliography = "（PPT 內未找到參考書目）"
    date_str: str | None = None
    ppt_r2_key: str | None = None
    if ppt_idx and ppt_idx <= len(ppt_files):
        ppt_path = ppt_files[ppt_idx - 1]
        print(f"  PPT [{ppt_idx}]: {ppt_path.name}")
        try:
            ppt_text = load_ppt_text(ppt_path)
            bibliography = load_ppt_bibliography(ppt_path)
        except Exception as e:
            print(f"  PPT read failed: {e}")
        date_str = ppt_path.name[:10].replace(".", "-")
        ppt_r2_key = f"qiangmian-ppt/ep{ppt_idx:02d}.pptx"

    if not date_str:
        d = entry.get("upload_date", "")
        date_str = f"{d[:4]}-{d[4:6]}-{d[6:]}" if len(d) == 8 else None

    # 取得 raw transcript：Whisper 或 DB
    if do_transcribe:
        audio = download_audio(video_id, label)
        if not audio:
            print(f"  ❌ audio download failed")
            return False
        try:
            raw = transcribe_audio_whisper(audio)
            print(f"  Whisper output: {len(raw):,} chars")
            if len(raw) < 1000:
                print(f"  ⚠ suspiciously short, abort ep {ep}")
                return False
        except Exception as e:
            print(f"  ❌ whisper error: {e}")
            return False
    else:
        row = fetch_row(ep)
        if not row:
            print(f"  ❌ no DB row")
            return False
        raw = extract_polished_raw(row["content"])
        print(f"  DB raw: {len(raw):,} chars")

    # Haiku 潤稿
    print("  Haiku polishing...")
    try:
        polished = polish_with_haiku(client, raw, raw_title, ppt_text)
        print(f"  Polished: {len(polished):,} chars")
    except Exception as e:
        print(f"  ❌ polish failed: {e}")
        return False

    # 驗證 sanity
    if len(polished) < 1500:
        print(f"  ⚠ polished too short, abort ep {ep}")
        return False
    if "##" not in polished:
        print(f"  ⚠ no ## headings detected, abort ep {ep}")
        return False

    # 組合最終內容
    final = build_final_content(raw_title, ep, date_str, polished, bibliography)
    ok = upsert_row(ep, raw_title, final, date_str, video_id, ppt_r2_key)
    print(f"  upsert: {'✓' if ok else '❌'}")

    if ok and do_transcribe:
        cleanup_audio(label)
    return ok


def main():
    print("=" * 60)
    print("千面上帝 整晚自動化")
    print(f"Phase 1 (polish only): {POLISH_ONLY}")
    print(f"Phase 2 (transcribe + polish): {TRANSCRIBE}")
    print("=" * 60)

    # Lazy installs
    for pkg in ("yt-dlp", "faster-whisper", "anthropic", "python-pptx", "requests"):
        try: __import__(pkg.replace("-", "_"))
        except ImportError:
            print(f"Installing {pkg}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pkg, "-q"])

    print("\nFetching playlist...")
    entries = fetch_playlist()
    print(f"  {len(entries)} videos in playlist")

    ppt_files = sorted_ppt_files()
    print(f"  {len(ppt_files)} PPT files")

    client = make_anthropic_client()
    print(f"\nUsing Haiku model: {HAIKU_MODEL}")

    # Phase 1: polish only
    print("\n" + "=" * 60)
    print(f"PHASE 1: polishing {len(POLISH_ONLY)} episodes")
    print("=" * 60)
    p1_ok = 0
    for ep in POLISH_ONLY:
        try:
            if process_episode(client, ep, entries, ppt_files, do_transcribe=False):
                p1_ok += 1
        except Exception as e:
            print(f"  fatal error on ep {ep}: {e}")
        time.sleep(2)
    print(f"\nPhase 1 done: {p1_ok}/{len(POLISH_ONLY)} OK")

    # Phase 2: transcribe + polish (slow)
    print("\n" + "=" * 60)
    print(f"PHASE 2: transcribing + polishing {len(TRANSCRIBE)} episodes")
    print("=" * 60)
    p2_ok = 0
    for ep in TRANSCRIBE:
        try:
            if process_episode(client, ep, entries, ppt_files, do_transcribe=True):
                p2_ok += 1
        except Exception as e:
            print(f"  fatal error on ep {ep}: {e}")
        time.sleep(2)
    print(f"\nPhase 2 done: {p2_ok}/{len(TRANSCRIBE)} OK")

    print("\n" + "=" * 60)
    print(f"ALL DONE. Phase 1: {p1_ok}/{len(POLISH_ONLY)}, Phase 2: {p2_ok}/{len(TRANSCRIBE)}")
    print("=" * 60)


if __name__ == "__main__":
    main()
